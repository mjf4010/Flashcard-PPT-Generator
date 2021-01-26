# Import the necessary modules
from fpdf import FPDF
import pandas as pd
import numpy as np
import requests
import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
# Define a function to delete slides
def delete_slide(presentation, index):
        xml_slides = presentation.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[index])      
# Create a function place the image in the center
def _add_image(slide, placeholder_id, image_path):
    placeholder = slide.placeholders[placeholder_id]
    # Calculate the image size of the image
    im = Image.open(image_path)
    width, height = im.size
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width
    # Insert the picture
    placeholder = placeholder.insert_picture(image_path)
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side
# Read the XLXS file into a DataFrame
data = pd.read_excel("input.xlsx", header = None)
# Name the columns so they can be manipulated
data.columns = ['URL', 'Word', 'Tip']
# Extract the Folder Name
name = data['Word'][0]
# Check if name is NAN
if pd.isnull(name):
    name = 'NO NAME WHYYYYYYYYY'
else:
    pass
# Drop the first 2 rows
data = data[2:].reset_index(drop = True)
# Create a list of URLs so we can download the images
urlList = data['URL'].to_list()
# Create a list of the tags
wordsList = data['Word'].to_list()
# Get the current working directory
path = os.getcwd()
# Create a folder for this PPTX
if not os.path.exists('{}'.format(name)):
    os.mkdir('{}'.format(name))
else:
    pass
# Create an Images directory if one doesn't already exist
if not os.path.exists('{}\{}\Images'.format(path, name)):
    os.mkdir('{}\{}\Images'.format(path, name))
else:
    pass
# Download the images from the URLs list
image = {}
n = 0
for url in urlList:
    n += 1
    response = requests.get(url)
    image[n] = open('{}\Images\{}.png'.format(name, n), 'wb')
    image[n].write(response.content)
    image[n].close()
# Instantiate a presentation object
prs = Presentation('template.pptx')
# Add the Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = name
subtitle.text = "Generated with Michael Tr.'s Python Script"
# Add the slides and images to the slides
layout = prs.slide_layouts[9]
slide = {}
for i in image:
    slide[i] = prs.slides.add_slide(layout)
    _add_image(slide[i], 13, '{}\Images\{}.png'.format(name, i))
# Add the words to the slides
x = 0
tit = {}
font = {}
sub = {}
for i in wordsList:
    x += 1
    text_frame = slide[x].placeholders[0].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = i
    p.alignment = PP_ALIGN.CENTER
    font = run.font
    font.size = Pt(72)
    font.bold = True
# Delete the first slide
delete_slide(prs, 0)
# Save the pptx file
prs.save('{}\{}.pptx'.format(name, name))
# Start the pptx file that was just created
os.startfile('{}\{}.pptx'.format(name, name))